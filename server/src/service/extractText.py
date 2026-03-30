import io
import csv
import fitz                          # PyMuPDF
from docx import Document
from pptx import Presentation
import openpyxl
from bs4 import BeautifulSoup
from faster_whisper import WhisperModel


def extract_pdf(data: bytes) -> str:
    with fitz.open(stream=data, filetype="pdf") as doc:
        return "\n".join(page.get_text() for page in doc)


def extract_docx(data: bytes) -> str:
    doc = Document(io.BytesIO(data))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def extract_pptx(data: bytes) -> str:
    prs = Presentation(io.BytesIO(data))
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text)
    return "\n".join(lines)


def extract_xlsx(data: bytes) -> str:
    wb = openpyxl.load_workbook(
        io.BytesIO(data), read_only=True, data_only=True
    )
    lines = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            line = "\t".join(str(c) for c in row if c is not None)
            if line.strip():
                lines.append(line)
    return "\n".join(lines)


def extract_csv(data: bytes) -> str:
    text = data.decode("utf-8", errors="ignore")
    reader = csv.reader(io.StringIO(text))
    return "\n".join(", ".join(row) for row in reader)


def extract_txt(data: bytes) -> str:
    return data.decode("utf-8", errors="ignore")


def extract_html(data: bytes) -> str:
    soup = BeautifulSoup(data, "lxml")
    return soup.get_text(separator="\n", strip=True) 


whisper_model = WhisperModel("tiny", device="cpu", compute_type="int8")
def extract_audio(data: bytes) -> str:
    import tempfile, os
    suffix = ".wav"  # whisper works best with wav
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        tmp.write(data)
        tmp_path = tmp.name

    try:
        segments, _ = whisper_model.transcribe(tmp_path, beam_size=1)
        return " ".join(segment.text for segment in segments)
    finally:
        os.remove(tmp_path)